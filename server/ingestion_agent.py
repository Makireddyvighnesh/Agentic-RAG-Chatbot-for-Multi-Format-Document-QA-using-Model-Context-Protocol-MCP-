import os
import json
from typing import List, Dict, Any

import pandas as pd
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

from mcp.server.fastmcp import FastMCP

# --- Configuration ---
CHUNK_DB_PATH = os.path.join(os.path.dirname(__file__), '..', 'data', 'chunks.json')
os.makedirs(os.path.dirname(CHUNK_DB_PATH), exist_ok=True)

# --- Document Parsers ---
def parse_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    return "".join(page.extract_text() for page in reader.pages)

def parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs)

def parse_pptx(file_path: str) -> str:
    pres = Presentation(file_path)
    return "\n".join(shape.text for slide in pres.slides for shape in slide.shapes if shape.has_text_frame)

def parse_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    return df.to_string()

def parse_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

# --- Agent Logic ---
mcp = FastMCP("ingestion-agent")

@mcp.tool()
def process_files(file_paths: List[str]) -> Dict[str, Any]:
    """
    Parses various document formats, splits them into text chunks,
    and saves them to a central data store.

    Args:
        file_paths: A list of absolute or relative paths to the documents.
    """
    print("Calling process_files")
    print(f"[IngestionAgent] Received request to process: {file_paths}", flush=True)
    
    full_text = ""
    for path in file_paths:
        ext = os.path.splitext(path)[1].lower()
        print(f"[IngestionAgent] Parsing {path}...", flush=True)
        try:
            if ext == '.pdf':
                full_text += parse_pdf(path)
            elif ext == '.docx':
                full_text += parse_docx(path)
            elif ext == '.pptx':
                full_text += parse_pptx(path)
            elif ext == '.csv':
                full_text += parse_csv(path)
            elif ext in ['.txt', '.md']:
                full_text += parse_txt(path)
            else:
                print(f"[IngestionAgent] WARNING: Unsupported file type: {ext}", flush=True)
                continue
            full_text += "\n\n" # Add separator between documents
        except Exception as e:
            print(f"[IngestionAgent] ERROR: Failed to parse {path}: {e}", flush=True)
            return {"status": "error", "message": f"Failed to parse {path}: {e}"}

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text=full_text)
    
    print(f"[IngestionAgent] Created {len(chunks)} chunks. Saving to {CHUNK_DB_PATH}", flush=True)
    with open(CHUNK_DB_PATH, 'w', encoding='utf-8') as f:
        json.dump(chunks, f)

    return {"status": "success", "chunks_created": len(chunks)}

if __name__ == "__main__":
    print("[IngestionAgent] Starting up...", flush=True)
    mcp.run(transport='stdio')

